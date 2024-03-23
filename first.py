from PyQt6.QtWidgets import QMessageBox, QVBoxLayout, QHBoxLayout, QListView, QLabel, QApplication, QWidget, QMainWindow, QPushButton, QFileDialog
from PyQt6.QtGui import QImage, QPixmap, QIcon
from PyQt6 import QtCore, uic
from PyQt6.QtCore import QObject, pyqtSignal, QThread
import cv2
# Only needed for access to command line arguments
import sys
from utils.detector import Detector
from pptx import Presentation
from pptx.util import Inches
import os

class First(QWidget):
    def __init__(self):
        super(First, self).__init__()
        # ui_file_path = os.path.join(sys._MEIPASS, 'first_page.ui')
        uic.loadUi('first_page.ui', self)
        

        # # Set the central widget of the Window.
        # self.setCentralWidget(self.uploadButton)
        self.LoadButton.clicked.connect(lambda: self.open("load"))
        self.ChooseOutputButton.clicked.connect(lambda: self.open("output"))
        self.detection_start.clicked.connect(self.start_detection)
        self.SaveButton.clicked.connect(self.save_image)
        self.filename = ""
        self.output = ""
        self.scroll_widget = QWidget()
        self.scroll_widget.setLayout(self.horizontalLayout)
        self.ImageScrollArea.setWidget(self.scroll_widget)
        self.ImageScrollArea.setWidgetResizable(True)
        self.detected_frames = []
        self.progressBar.setValue(0)
        self.threshold_slider.setValue(99)
        self.threshold_slider.valueChanged.connect(self.threshold_display)

        self.worker = None

        # main_layout.addWidget(self.ImageScrollArea)
        # main_layout.addWidget(self.progressBar)
        # main_layout.addWidget(self.threshold_slider)
        # main_layout.addWidget(self.SaveButton)
        # main_layout.addWidget(self.filenametext)
        # main_layout.addWidget(self.output_text)
        # main_layout.addWidget(self.threshold_text)
        # main_layout.addWidget(self.ImageShowLabel)
        # main_layout.addWidget(self.youtubepage_button)
        # self.button_layout.setParent(None)
        

    def open(self, method):
        if method == "load":
            filename = QFileDialog.getOpenFileName(self, 'Open File', '.')
            self.filename = filename[0]
            if self.filename:
                if not self.filename.split(".")[-1] == "mp4" and not self.filename.split(".")[-1] == "avi":
                    msg = QMessageBox()
                    msg.setIcon(QMessageBox.Icon.Critical)
                    msg.setText("Error")
                    msg.setInformativeText('Wrong file type (Only allow .mp4 and .avi)')
                    msg.setWindowTitle("Error")
                    msg.exec()
                else:
                    self.read_first_image()
                    self.filenametext.setText("Uploaded File: " + filename[0])

        elif method == "output":
            filename = QFileDialog.getExistingDirectory(self, 'Open File', '.')
            self.output_text.setText("Output directory: " + filename)
            self.output = filename

    def start_detection(self):
        if self.filename and self.output:
            self.clear()
            self.worker = DetectorWorker(self.filename, threshold=self.threshold_slider.value() / 100)
            self.thread = QThread()
            self.worker.moveToThread(self.thread)
            self.worker.progressChanged.connect(self.track_progress)
            self.worker.detectionFinished.connect(self.create_slider_buttons)
            self.worker.finished.connect(self.thread.quit)
            self.worker.finished.connect(self.count_images)

            self.thread.started.connect(self.worker.detect_slides)
            self.thread.start()


    def create_slider_buttons(self, detected_frames):
        self.progressBar.setValue(100)
        self.detected_frames = detected_frames
        for i, frame in enumerate(self.detected_frames):
            trainsient_button = SlideButton(frame.img)
            trainsient_button.clicked.connect(lambda state, idx=i: self.show_image(self.detected_frames[idx].img, idx))
            self.horizontalLayout.addWidget(trainsient_button)
            
    def read_first_image(self):
        cap = cv2.VideoCapture(self.filename)
        ret, frame = cap.read()
        self.show_image(frame)
        

    def show_image(self, img, idx=None):
        img = cv2.resize(img, (self.ImageShowLabel.width(), self.ImageShowLabel.height()))
        h, w, d = img.shape
        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        previewImage = QImage(img.data, w, h, d * w, QImage.Format.Format_RGB888)
        self.ImageShowLabel.setPixmap(QPixmap.fromImage(previewImage))

        if idx != None:
            self.detected_frames[idx].marked = not self.detected_frames[idx].marked
        
        self.count_images()

    def save_image(self):
        new_folder = self.filename.split("/")[-1].split(".")[0]
        new_path = self.output + "/" + new_folder
        if not os.path.exists(new_path):
            os.makedirs(new_path)
        if self.detected_frames:
            for img in self.detected_frames:
                if img.marked:
                    cv2.imwrite(new_path + "/" + img.time + ".jpg", img.img)
        self.create_ppt([new_path + "/" + img.time + ".jpg" for img in self.detected_frames if img.marked], new_path + "/presentation.pptx")
    
    def count_images(self):
        count = 0
        for frame in self.detected_frames:
            if frame.marked:
                count += 1
        self.slides_selected.setText(f"Slides selected: {count}")

    def threshold_display(self):
        self.threshold_text.setText("Threshold: " + str(self.threshold_slider.value() / 100))
    
    def track_progress(self, value):
        self.progressBar.setValue(round(value))

    def create_ppt(self, images, output):
        prs = Presentation()

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = top = Inches(1)
            pic = slide.shapes.add_picture(img, left, top, width=Inches(8), height=Inches(6))

        prs.save(output)

    def clear(self):
        self.detected_frames = []
        for i in reversed(range(self.horizontalLayout.count())): 
            self.horizontalLayout.itemAt(i).widget().setParent(None)

class DetectorWorker(QObject):
    progressChanged = pyqtSignal(int)
    detectionFinished = pyqtSignal(list)
    finished = pyqtSignal()
    
    def __init__(self, filename, threshold=0.99):
        super(DetectorWorker, self).__init__()
        self.filename = filename
        self.threshold = threshold
    
    def detect_slides(self):
        detector = Detector(self.filename, "output", ".jpg", threshold=self.threshold)
        detected_frames = detector.detect_slides(self.track_progress)
        self.detectionFinished.emit(detected_frames)
        self.finished.emit()
    
    def track_progress(self, value):
        self.progressChanged.emit(int(value))

class SlideButton(QPushButton):
    def __init__(self, img):
        super(SlideButton, self).__init__()
        self.img = img
        frame = cv2.resize(self.img, (100, 100))
        h, w, d = frame.shape
        frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)

        transient_image = QImage(frame.data, w, h, d * w, QImage.Format.Format_RGB888)

        self.setFixedSize(QtCore.QSize(100, 100))
        self.setIcon(QIcon(QPixmap.fromImage(transient_image)))
        self.setIconSize(QtCore.QSize(75, 75))
        self.setCheckable(True)
        self.setChecked(True)


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = First()
    window.show()
    sys.exit(app.exec_())